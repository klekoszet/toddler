import pandas as pd
import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class ExportManager:
    def __init__(self, df: pd.DataFrame, history: list, prep_summary: dict, clean_opts: dict, scale_items: list, time_col: str, speeder_threshold: float):
        self.df = df
        self.history = history 
        self.prep_summary = prep_summary
        self.clean_opts = clean_opts or {}
        self.scale_items = scale_items or []
        self.time_col = time_col
        self.speeder_threshold = speeder_threshold
        self.filepath = f"Raport_Wundt_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"

    def export_excel(self):
        if not self.history: return
        path = self.filepath + ".xlsx"
        try:
            with pd.ExcelWriter(path, engine='openpyxl') as writer:
                if self.prep_summary:
                    pd.DataFrame(list(self.prep_summary.items())).to_excel(writer, sheet_name="Data_Prep", index=False)
                for i, batch in enumerate(self.history):
                    if not batch: continue
                    df = pd.DataFrame(batch) if isinstance(batch, list) else pd.DataFrame([batch])
                    title = batch[0].get("Wybrany Test", f"Analiza_{i+1}") if isinstance(batch, list) else "Analiza"
                    df.to_excel(writer, sheet_name=f"{i+1}_{title}"[:30].replace("/", ""), index=False)
            print(f"[SUKCES] Raport Excel został wygenerowany: {path}")
        except Exception as e: 
            print(f"[Błąd] Excel: {e}")
            
    def export_pptx(self):
        if not PPTX_AVAILABLE or not self.history: return
        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Raport Analityczny WUNDT"
            for i, batch in enumerate(self.history):
                if not batch or not isinstance(batch, list): continue
                slide = prs.slides.add_slide(prs.slide_layouts[5]) 
                slide.shapes.title.text = f"{batch[0].get('Wybrany Test', 'Analiza')} (#{i+1})"
                rows, cols = len(batch) + 1, len(batch[0].keys())
                table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * rows)).table
                headers = list(batch[0].keys())
                for col_idx, header in enumerate(headers): table.cell(0, col_idx).text = str(header)
                for row_idx, res in enumerate(batch):
                    for col_idx, key in enumerate(headers): table.cell(row_idx + 1, col_idx).text = str(res.get(key, ""))
            prs.save(self.filepath + ".pptx")
            print(f"[SUKCES] Prezentacja PPTX została wygenerowana: {self.filepath}.pptx")
        except Exception as e: 
            print(f"[Błąd] PPTX: {e}")

    def _fmt_spss(self, val):
        if pd.isna(val) or val == "": return "1"
        if isinstance(val, (int, float)): return str(int(val)) if float(val).is_integer() else str(val)
        val_str = str(val).strip()
        try:
            num = float(val_str)
            return str(int(num)) if num.is_integer() else str(num)
        except ValueError:
            return f"'{val_str}'"

    def export_spss_syntax(self, global_mode: str):
        path = self.filepath + "_syntax.sps"
        lines = ["* =========================================.", "* SYNTAX WYGENEROWANY AUTOMATYCZNIE - WUNDT.", "* =========================================.\n"]
        
        lines.append("* --- ETAP 1: DATA PREP ---.")
        if global_mode == '1': 
            if self.clean_opts.get('missing') and self.scale_items:
                sc_str = " ".join(self.scale_items)
                lines.append(f"COUNT Missing_Count = {sc_str} (SYSMIS).\nSELECT IF (Missing_Count <= {int(len(self.scale_items) * 0.15)}).\nEXECUTE.")
            if self.clean_opts.get('straight_liners') and self.scale_items:
                lines.append(f"COMPUTE Var_Scale = VARIANCE({' TO '.join([self.scale_items[0], self.scale_items[-1]])}).\nSELECT IF (Var_Scale >= 0.1).\nEXECUTE.")
            if self.clean_opts.get('speeders') and self.time_col and self.speeder_threshold:
                lines.append(f"SELECT IF ({self.time_col} >= {self.speeder_threshold}).\nEXECUTE.")
        else: lines.append("* Baza zostala w pelni oczyszczona i zwazona przez Pythona (Zapisano plik .SAV).")

        lines.append("\n* --- ETAP 2: ANALIZY STATYSTYCZNE ---.")
        for i, batch in enumerate(self.history):
            if not isinstance(batch, list) or len(batch) == 0: continue
            lines.append(f"\n* Analiza #{i+1}.")
            
            for row in batch:
                test = str(row.get("Wybrany Test", "")).lower()
                v1 = str(row.get("Zmienna Zależna", row.get("Zmienna 1", row.get("Zmienna", ""))))
                v2 = str(row.get("Zmienna Grupująca", row.get("Zmienna 2", "")))
                if not v1: continue

                if "t2b" in test:
                    lines.append(f"RECODE {v1} (4 THRU 5 = 100) (ELSE = 0) INTO {v1}_T2B.\nDESCRIPTIVES VARIABLES={v1}_T2B /STATISTICS=MEAN.")
                elif "nps" in test:
                    lines.append(f"RECODE {v1} (9 THRU 10 = 100) (7 THRU 8 = 0) (0 THRU 6 = -100) INTO {v1}_NPS.\nDESCRIPTIVES VARIABLES={v1}_NPS /STATISTICS=MEAN.")
                elif "indeks" in test:
                    lines.append(f"COMPUTE {row.get('Utworzony Indeks', 'Indeks')} = MEAN({row.get('Składowe', '').replace(',', '')}).")
                elif "korelacja" in test:
                    lines.append(f"CORRELATIONS /VARIABLES={v1} {v2}.")
                elif "zależnych" in test:
                    lines.append(f"T-TEST PAIRS={v1} WITH {v2} (PAIRED).")
                elif "wilcoxona" in test:
                    lines.append(f"NPAR TESTS /WILCOXON={v1} WITH {v2} (PAIRED).")
                elif "t-studenta" in test:
                    uv = sorted(self.df[v2].dropna().unique().tolist()) if v2 in self.df.columns else []
                    val1 = self._fmt_spss(uv[0]) if len(uv) > 0 else "1"
                    val2 = self._fmt_spss(uv[1]) if len(uv) > 1 else "2"
                    lines.append(f"T-TEST GROUPS={v2}({val1} {val2}) /VARIABLES={v1}.")
                elif "manna-whitneya" in test:
                    uv = sorted(self.df[v2].dropna().unique().tolist()) if v2 in self.df.columns else []
                    val1 = self._fmt_spss(uv[0]) if len(uv) > 0 else "1"
                    val2 = self._fmt_spss(uv[1]) if len(uv) > 1 else "2"
                    lines.append(f"NPAR TESTS /M-W={v1} BY {v2}({val1} {val2}).")
                elif "anova" in test:
                    lines.append(f"ONEWAY {v1} BY {v2} /STATISTICS DESCRIPTIVES HOMOGENEITY WELCH.")
                elif "kruskala" in test:
                    min_v = self._fmt_spss(self.df[v2].min()) if v2 in self.df.columns else "1"
                    max_v = self._fmt_spss(self.df[v2].max()) if v2 in self.df.columns else "3"
                    lines.append(f"NPAR TESTS /K-W={v1} BY {v2}({min_v} {max_v}).")
            lines.append("EXECUTE.")

        with open(path, "w", encoding="utf-8-sig") as f: f.write("\n".join(lines))
        print(f"[Eksport OK] Zapisano plik SPSS Syntax: {path}")